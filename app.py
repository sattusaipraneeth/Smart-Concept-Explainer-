# ✅ EduGenius AI (Professional UI)

import streamlit as st
import fitz  # PyMuPDF
import google.generativeai as genai
from gtts import gTTS
from pptx import Presentation
import tempfile
import re
import os
import io
import time
from fpdf import FPDF
from pydub import AudioSegment
import base64


# === PAGE SETUP === (MOVED TO TOP)
st.set_page_config(
    page_title="EduGenius AI",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded"
)

# === CONFIG ===
GEMINI_API_KEY = "AIzaSyA60Bzk8pQc50tq8voxQpdh6AsZIyzbRk8"
genai.configure(api_key=GEMINI_API_KEY)
gemini_model = genai.GenerativeModel("gemini-1.5-flash")

# === CUSTOM CSS ===
def inject_custom_css():
    st.markdown("""
    <style>
        html, body, .main {
            background-color: var(--background-color);
            color: var(--text-color);
        }

        /* Headers */
        h1, h2, h3, h4, h5, h6 {
            color: var(--text-color);
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        }

        /* Sidebar */
        .sidebar .sidebar-content {
            background-color: var(--background-color);
            box-shadow: 2px 0 10px rgba(0,0,0,0.1);
        }

        /* Buttons */
        .stButton>button {
            background-color: #4a6fa5;
            color: white;
            border-radius: 8px;
            padding: 10px 20px;
            font-weight: 600;
            transition: all 0.3s;
            border: none;
            box-shadow: 0 2px 5px rgba(0,0,0,0.1);
        }
        .stButton>button:hover {
            background-color: #3a5a80;
            transform: translateY(-2px);
            box-shadow: 0 4px 8px rgba(0,0,0,0.2);
        }

        /* Select boxes */
        .stSelectbox>div>div>select {
            border: 1px solid #dfe6e9;
            border-radius: 8px;
            padding: 10px;
            font-size: 16px;
        }

        /* File uploader */
        .stFileUploader>div>div {
            border: 2px dashed #bdc3c7;
            border-radius: 12px;
            padding: 25px;
            background-color: rgba(255,255,255,0.05);
        }

        /* Cards */
        .card {
            background: var(--secondary-background-color);
            border-radius: 12px;
            box-shadow: 0 6px 12px rgba(0,0,0,0.1);
            padding: 25px;
            margin-bottom: 25px;
            color: var(--text-color);
            border-left: 5px solid #4a6fa5;
        }

        /* Divider */
        .divider {
            border-top: 2px solid #e0e0e0;
            margin: 25px 0;
            opacity: 0.3;
        }

        /* Audio player description */
        .audio-caption {
            font-size: 14px;
            text-align: center;
            color: var(--text-color);
            margin-top: -10px;
        }
        
        /* New feature: Tag styling */
        .tag {
            display: inline-block;
            background-color: #4a6fa5;
            color: white;
            padding: 4px 8px;
            border-radius: 4px;
            font-size: 12px;
            margin-right: 5px;
            margin-bottom: 5px;
        }
        
        /* Progress bar styling */
        .stProgress > div > div > div {
            background-color: #4a6fa5;
        }
        
        /* Chat message styling */
        .stChatMessage {
            border-radius: 12px;
            padding: 15px;
            margin-bottom: 15px;
        }
    </style>
    <script>
    const observer = new MutationObserver((mutations, obs) => {
        document.documentElement.style.setProperty('--background-color', getComputedStyle(document.body).backgroundColor);
        document.documentElement.style.setProperty('--secondary-background-color', getComputedStyle(document.body).backgroundColor === 'rgb(255, 255, 255)' ? '#f8f9fa' : '#1e1e1e');
        document.documentElement.style.setProperty('--text-color', getComputedStyle(document.body).color);
    });
    observer.observe(document.body, { attributes: true, childList: true, subtree: true });
    </script>
    """, unsafe_allow_html=True)

inject_custom_css()

# === SESSION STATE ===
if 'concepts' not in st.session_state:
    st.session_state.concepts = []
if 'explanations' not in st.session_state:
    st.session_state.explanations = {}
if 'audio_files' not in st.session_state:
    st.session_state.audio_files = {}
if 'pdf_text' not in st.session_state:
    st.session_state.pdf_text = ""
if 'current_topic' not in st.session_state:
    st.session_state.current_topic = None
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'quiz_questions' not in st.session_state:
    st.session_state.quiz_questions = {}
if 'summaries' not in st.session_state:
    st.session_state.summaries = {}

# === HEADER ===
with st.container():
    col1, col2 = st.columns([1, 3])
    with col1:
        st.image("https://cdn-icons-png.flaticon.com/512/2232/2232688.png", width=120)
    with col2:
        st.markdown("""
        <h1 style='margin-bottom: 0;'>EduGenius AI</h1>
        <p style='color: #7f8c8d; margin-top: 0;'>AI-Powered Learning Companion with Multimodal Explanations</p>
        """, unsafe_allow_html=True)
    
    st.markdown("---")

# === SIDEBAR ===
with st.sidebar:
    st.markdown("""
    <div style='text-align: center; margin-bottom: 30px;'>
        <h3 style='color: #2c3e50;'>📂 Document Upload</h3>
    </div>
    """, unsafe_allow_html=True)
    
    uploaded_file = st.file_uploader(
        "Upload your lecture notes (PDF/PPTX)",
        type=["pdf", "pptx"],
        label_visibility="collapsed"
    )
    
    st.markdown("---")
    
    st.markdown("""
    <div style='text-align: center;'>
        <h4 style='color: #2c3e50;'>⚙️ Settings</h4>
    </div>
    """, unsafe_allow_html=True)
    
    selected_lang = st.selectbox(
        "Explanation Language",
        ["English", "Telugu", "Hindi", "Spanish", "French"],
        index=0,
        key="lang_select"
    )
    
    difficulty_level = st.selectbox(
        "Explanation Difficulty",
        ["Beginner", "Intermediate", "Advanced"],
        index=1,
        key="difficulty_select"
    )
    
    explanation_style = st.selectbox(
        "Explanation Style",
        ["Academic", "Conversational", "Storytelling", "Visual"],
        index=1,
        key="style_select"
    )

# === TEXT EXTRACTORS ===
def extract_text_from_pdf(uploaded_file):
    with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
        return "\n".join(page.get_text() for page in doc)

def extract_text_from_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

# === CLEAN FOR AUDIO ===
def clean_text_for_tts(text, lang):
    """Enhanced text cleaning for better TTS output"""
    # Common replacements for all languages
    replacements = {
        "-": "",
        "(": ", ",
        ")": ", ",
        "[": ", ",
        "]": ", ",
        "{": ", ",
        "}": ", ",
        "•": "",
        "\n": ". ",
        "\t": " ",
    }
    
    # Language-specific replacements
    if lang == "Telugu":
        lang_replacements = {
            "ఉదా.": "ఉదాహరణకు",
            "అంటే": "అంటే ",
            "=": " సమానం ",
            "+": " ప్లస్ ",
            "/": " డివైడెడ్బై ",
            "Dr.": "డాక్టర్ ",
            "e.g.": "ఉదాహరణకు ",
        }
    elif lang == "Hindi":
        lang_replacements = {
            "डॉ.": "डाक्टर ",
            "=": " बराबर ",
            "+": " प्लस ",
            "/": " भागा ",
            "उदा.": "उदाहरण के लिए ",
        }
    else:  # English and others
        lang_replacements = {
            "e.g.": "for example",
            "i.e.": "that is",
            "Dr.": "Doctor",
            "=": " equals ",
            "+": " plus ",
            "/": " divided by ",
        }
    
    replacements.update(lang_replacements)
    
    for k, v in replacements.items():
        text = text.replace(k, v)
    
    # Add pauses between sections
    text = re.sub(r"(\d+)\.", r"\n\n\1.", text)
    text = re.sub(r" +", " ", text)
    return text.strip()

# === IDENTIFY CONCEPTS ===
def identify_concepts(text):
    """
    Smart concept extractor that identifies main topics and subtopics from any type of academic PDF.
    Works across math, science, and general texts.
    """
    prompt = f"""
You are a highly accurate document parser. Analyze the following academic content and extract only meaningful, high-level topics and subtopics.

STRICT INSTRUCTIONS:
- Only extract headings, chapter names, or subheadings — not formulas, not examples.
- Avoid garbage like short codes (e.g., TP, OM) or numeric-only lines.
- Do not return duplicates or repeated wording (e.g., "Permutation" and "Permutations").
- Do not include the chapter summary or index blocks (like Unit III...).
- Preserve the logical order from the text.

Return only a clean numbered list of the concepts found in the below text:

{text[:10000]}

FORMAT STRICTLY LIKE THIS:
1. Topic One
2. Topic Two
3. ...
"""

    try:
        response = gemini_model.generate_content(prompt)
        lines = response.text.strip().split("\n")

        clean = []
        for line in lines:
            if ". " in line:
                content = line.split(". ", 1)[1].strip()
                if len(content) > 5 and not re.fullmatch(r"[0-9. ]+", content):
                    clean.append(content)

        return clean

    except Exception as e:
        st.error(f"❌ Gemini extraction error: {e}")
        return []

# === EXPLAIN CONCEPT ===
def explain_concept(concept, context, lang, difficulty="Intermediate", style="Conversational"):
    if lang != "English":
        prompt = f"""
Explain the topic '{concept}' from the given PDF in a **simple, clear and student-friendly** manner using **{lang} + English mix** style.

Your explanation must follow this structure:

1. **Definition**  
   - Concept in simple terms in {lang}  
   - Important terminology in English  
   - Example for clarity

2. **Key Characteristics**  
   - 5-6 important features  
   - Each point in simple sentences  
   - Keywords in English

3. **Real-life Examples**  
   - Examples from the PDF  
   - Additional examples (3 total)  
   - {lang} explanation + English terms

4. **Problems/Applications**  
   - Step-by-step problem solving  
   - Formulas in English, steps in {lang}  
   - Final answer

5. **Importance/Uses**  
   - Practical applications  
   - English terms explained

6. **Common Mistakes**  
   - Common student errors  
   - Examples to explain

Difficulty Level: {difficulty}
Style: {style}

🧠 Style:
- Use **spoken conversational {lang} + English mix**
- Very simple language
- Highlight important English terms
- No lengthy or complex sentences
- Make it feel like a teacher is explaining to a student in class

Context:  
{context[:10000]}
        """
    else:
        prompt = f"""
Explain the topic '{concept}' from the given PDF in clear and detailed English.

Follow this structure:
1. Definition
2. Key Characteristics
3. Real-life Examples
4. Problems/Applications
5. Importance/Uses
6. Common Mistakes

Difficulty Level: {difficulty}
Style: {style}

- Use {style.lower()} teaching style
- Each section should have appropriate detail for {difficulty} level
- Include formulae, diagrams, and examples if found in PDF

Context:  
{context[:10000]}
        """

    try:
        response = gemini_model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        st.error(f"Explanation generation failed: {str(e)}")
        return f"Error generating explanation. Please try again. {str(e)}"

# === GENERATE QUIZ ===
def generate_quiz_questions(concept, context, num_questions=5):
    prompt = f"""
Generate {num_questions} quiz questions about '{concept}' based on the following content.
Include multiple choice questions with 4 options each and mark the correct answer.

FORMAT STRICTLY LIKE THIS:
1. Question text?
   A) Option 1
   B) Option 2
   C) Option 3
   D) Option 4
   Answer: D

2. Next question...
   A) ...
   B) ...
   C) ...
   D) ...
   Answer: A

Context:
{context[:5000]}
"""
    try:
        response = gemini_model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        st.error(f"Quiz generation failed: {str(e)}")
        return None

# === GENERATE SUMMARY ===
def generate_summary(concept, context, lang="English"):
    prompt = f"""
Generate a concise summary of '{concept}' in {lang} based on the following content.
The summary should be no more than 150 words and highlight the key points.

Key requirements:
- Focus on main ideas only
- Use simple language
- Include key terms
- Keep it very brief

Context:
{context[:5000]}
"""
    try:
        response = gemini_model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        st.error(f"Summary generation failed: {str(e)}")
        return None

# === AUDIO FUNCTIONS ===
def generate_high_quality_audio(text, lang):
    """Generate high quality audio for multiple languages directly from text"""
    try:
        # Language code mapping
        lang_codes = {
            "English": "en",
            "Telugu": "te",
            "Hindi": "hi",
            "Spanish": "es",
            "French": "fr"
        }
        
        lang_code = lang_codes.get(lang, "en")
        
        # Clean text based on language
        cleaned_text = clean_text_for_tts(text, lang)
        
        # Create in-memory audio file
        audio_bytes = io.BytesIO()
        
        # Generate audio directly
        tts = gTTS(
            text=cleaned_text,
            lang=lang_code,
            slow=False,
            lang_check=False  # Bypass strict language checking
        )
        
        # Save to in-memory bytes
        tts.write_to_fp(audio_bytes)
        audio_bytes.seek(0)
        
        return audio_bytes
    
    except Exception as e:
        st.error(f"Audio generation failed: {str(e)}")
        return None

# === PDF DOWNLOAD ===
def download_explanations_as_pdf():
    import re
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Load and register Unicode font (required for multiple languages)
    font_path = "NotoSans-Regular.ttf"
    if not os.path.exists(font_path):
        # Fallback to built-in font if Noto font not found
        pdf.set_font("Arial", size=12)
    else:
        # Register all font variants we need
        pdf.add_font("Noto", "", font_path, uni=True)
        pdf.add_font("Noto", "B", font_path, uni=True)  # Same file for bold
        pdf.add_font("Noto", "I", font_path, uni=True)  # Same file for italic
        pdf.add_font("Noto", "BI", font_path, uni=True) # Same file for bold-italic
        pdf.set_font("Noto", size=12)

    # Add title
    pdf.set_font("Noto" if os.path.exists(font_path) else "Arial", size=16, style='B')
    pdf.cell(0, 10, txt="EduGenius AI - Concept Explanations", ln=True, align='C')
    pdf.ln(10)
    pdf.set_font("Noto" if os.path.exists(font_path) else "Arial", size=12)

    for concept, explanation in st.session_state.explanations.items():
        # Set title in bold
        pdf.set_font("Noto" if os.path.exists(font_path) else "Arial", size=14, style='B')
        pdf.cell(0, 10, txt=concept, ln=True)
        pdf.set_font("Noto" if os.path.exists(font_path) else "Arial", size=12)
        
        # Add summary if available
        if concept in st.session_state.summaries:
            pdf.set_font("Noto" if os.path.exists(font_path) else "Arial", style='I')
            pdf.multi_cell(0, 8, txt=f"Summary: {st.session_state.summaries[concept]}")
            pdf.ln(5)
            pdf.set_font("Noto" if os.path.exists(font_path) else "Arial", size=12)

        # Clean up markdown-style text
        cleaned = re.sub(r"\*\*(.*?)\*\*", r"\1", explanation)  # remove bold markdown
        cleaned = re.sub(r"##+", "", cleaned)  # remove headings like ## Title
        cleaned = re.sub(r"`", "", cleaned)    # remove backticks
        cleaned = re.sub(r"\s+", " ", cleaned) # normalize spaces
        cleaned = re.sub(r"•", "-", cleaned)   # replace bullets if needed

        # Add text to PDF
        pdf.multi_cell(0, 8, cleaned + "\n")

        # Add quiz questions if available
        if concept in st.session_state.quiz_questions:
            pdf.set_font("Noto" if os.path.exists(font_path) else "Arial", style='B')
            pdf.cell(0, 10, txt="Quiz Questions:", ln=True)
            pdf.set_font("Noto" if os.path.exists(font_path) else "Arial", size=12)
            
            questions = st.session_state.quiz_questions[concept].split("\n")
            for line in questions:
                if line.strip():
                    pdf.multi_cell(0, 8, txt=line)
            
            pdf.ln(5)

        # Divider line
        pdf.ln(5)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(10)

    # Create a temporary file to store the PDF
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    temp_path = temp_file.name
    pdf.output(temp_path)
    temp_file.close()
    
    # Read the PDF bytes
    with open(temp_path, "rb") as f:
        pdf_bytes = f.read()
    
    # Clean up the temporary file
    os.unlink(temp_path)
    
    return pdf_bytes

# === MAIN CONTENT ===
if uploaded_file and not st.session_state.get('pdf_text', ""):
    with st.spinner("📄 Extracting content from document..."):
        ext = uploaded_file.name.split(".")[-1].lower()
        if ext == "pdf":
            st.session_state.pdf_text = extract_text_from_pdf(uploaded_file)
        elif ext == "pptx":
            st.session_state.pdf_text = extract_text_from_pptx(uploaded_file)
    
    with st.spinner("🧠 Analyzing document for key concepts..."):
        st.session_state.concepts = identify_concepts(st.session_state.pdf_text)
        if st.session_state.concepts:
            st.success("✅ Document processed successfully!")
            with st.expander("📋 Extracted Concepts", expanded=True):
                cols = st.columns(2)
                for i, concept in enumerate(st.session_state.concepts):
                    with cols[i % 2]:
                        st.info(f"• {concept}")

# === CONCEPT EXPLANATION SECTION ===
if st.session_state.get('concepts', []):
    st.markdown("---")

    with st.container():
        st.markdown("""
        <div style='text-align: center; margin-bottom: 20px;'>
            <h2>Concept Explanation</h2>
            <p>Select a concept to explore in detail</p>
        </div>
        """, unsafe_allow_html=True)

        selected_topic = st.selectbox(
            "Select a concept to explain:",
            st.session_state.concepts,
            key="topic_select"
        )

        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("📝 Generate Explanation", key="explain_btn", use_container_width=True):
                with st.spinner(f"Generating {selected_lang} explanation..."):
                    explanation = explain_concept(
                        selected_topic,
                        st.session_state.pdf_text,
                        selected_lang,
                        difficulty_level,
                        explanation_style
                    )
                    st.session_state.explanations[selected_topic] = explanation
                    
                    # Generate summary automatically
                    summary = generate_summary(selected_topic, st.session_state.pdf_text, selected_lang)
                    if summary:
                        st.session_state.summaries[selected_topic] = summary

        with col2:
            if st.button("🎧 Generate Audio", key="audio_btn", use_container_width=True):
                if selected_topic in st.session_state.explanations:
                    with st.spinner("Generating audio explanation..."):
                        audio_bytes = generate_high_quality_audio(
                            st.session_state.explanations[selected_topic], 
                            selected_lang
                        )
                        if audio_bytes:
                            st.session_state.audio_files[selected_topic] = audio_bytes
                            st.success("Audio generated successfully!")
                            # Play the audio immediately
                            st.audio(audio_bytes, format='audio/mp3')
                            st.markdown("""
                            <div style="font-size: 14px; text-align: center; color: var(--text-color); margin-top: -10px;">
                                🎧 <em>Playback speed: 1.25x recommended for better clarity</em>
                            </div>
                            """, unsafe_allow_html=True)
                        else:
                            st.error("Failed to generate audio")
                else:
                    st.warning("Please generate the explanation first")

        with col3:
            if st.button("📝 Generate Quiz", key="quiz_btn", use_container_width=True):
                if selected_topic in st.session_state.explanations:
                    with st.spinner("Generating quiz questions..."):
                        quiz = generate_quiz_questions(selected_topic, st.session_state.pdf_text)
                        if quiz:
                            st.session_state.quiz_questions[selected_topic] = quiz
                            st.success("Quiz generated successfully!")
                else:
                    st.warning("Please generate the explanation first")

        if selected_topic in st.session_state.get('explanations', {}):
            st.markdown("---")
            with st.container():
                # Display summary if available
                if selected_topic in st.session_state.get('summaries', {}):
                    with st.expander("📌 Key Summary", expanded=True):
                        st.info(st.session_state.summaries[selected_topic])
                
                # Main explanation card
                st.markdown(f"""
                <div class='card'>
                    <h3 style='color: #4a6fa5;'>{selected_topic}</h3>
                    <div class='divider'></div>
                    {st.session_state.explanations[selected_topic]}
                </div>
                """, unsafe_allow_html=True)
                
                # Quiz section
                if selected_topic in st.session_state.get('quiz_questions', {}):
                    with st.expander("🧠 Quiz Questions", expanded=False):
                        st.markdown(st.session_state.quiz_questions[selected_topic])

# === EXPORT SECTION ===
if st.session_state.get('explanations', {}):
    st.markdown("---")
    with st.container():
        st.markdown("""
        <div style='text-align: center; margin-bottom: 20px;'>
            <h2>Export Options</h2>
            <p>Download your learning materials for offline use</p>
        </div>
        """, unsafe_allow_html=True)
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            if st.button("📄 Download PDF", key="pdf_btn", use_container_width=True):
                with st.spinner("Generating PDF..."):
                    pdf_bytes = download_explanations_as_pdf()
                    st.download_button(
                        label="⬇️ Download Now",
                        data=pdf_bytes,
                        file_name="EduGenius_Explanations.pdf",
                        mime="application/pdf",
                        key="pdf_download"
                    )
        
        with col2:
            # Export all audio as ZIP (placeholder)
            if st.button("🎧 Export All Audio", use_container_width=True):
                st.warning("Feature coming soon!")
        
        with col3:
            # Export as Anki deck (placeholder)
            if st.button("📚 Export as Anki Deck", use_container_width=True):
                st.warning("Feature coming soon!")

# === CHAT SECTION ===
st.markdown("---")
st.markdown("""
<div style='text-align: center; margin-bottom: 20px;'>
    <h2>🧠 EduGenius Assistant</h2>
    <p>Ask questions about any concept (Chat Memory Enabled)</p>
</div>
""", unsafe_allow_html=True)

# Display chat history
for chat in st.session_state.get('chat_history', []):
    with st.chat_message("user", avatar="🧑‍💻"):
        st.markdown(chat["user"])
    with st.chat_message("assistant", avatar="🤖"):
        st.markdown(chat["assistant"])

# Chat input box
user_question = st.chat_input("Type your question here...")
if user_question:
    with st.chat_message("user", avatar="🧑‍💻"):
        st.markdown(user_question)

    with st.chat_message("assistant", avatar="🤖"):
        chat_model = genai.GenerativeModel("gemini-1.5-flash")
        response = chat_model.generate_content(
            f"Context: {st.session_state.get('pdf_text', '')[:2000]}\n\n"
            f"Previous chat history: {st.session_state.get('chat_history', [])[-3:]}\n\n"
            f"Question: {user_question}\n\n"
            f"Answer in {selected_lang} using {explanation_style.lower()} style for a {difficulty_level.lower()} level learner."
        )
        st.markdown(response.text)

        # Save to chat history
        if 'chat_history' not in st.session_state:
            st.session_state.chat_history = []
        st.session_state.chat_history.append({
            "user": user_question,
            "assistant": response.text
        })

# === FOOTER ===
st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #7f8c8d; font-size: 14px;'>
    <p>EduGenius AI - Your Smart Learning Companion</p>
    <p>© 2023 EduGenius AI. All rights reserved.</p>
</div>
""", unsafe_allow_html=True)

